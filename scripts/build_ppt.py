"""
魔法大脑 - 方案 PPT 生成脚本
============================
用 python-pptx 生成 16:9 深色科技风幻灯片。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "screenshots"
OUT = ROOT / "docs" / "魔法大脑-方案PPT.pptx"

# 配色 (深空科技风, 跟产品UI一致)
BG = RGBColor(0x05, 0x07, 0x0F)
CARD = RGBColor(0x0F, 0x17, 0x2A)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)     # 青
ACCENT2 = RGBColor(0x22, 0xD3, 0xEE)    # 青2
SAVE = RGBColor(0x4A, 0xDE, 0x80)       # 绿(省钱)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
DIM = RGBColor(0x94, 0xA3, 0xB8)
FAINT = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return s


def add_rect(slide, left, top, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_text(slide, left, top, w, h, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT, font="Manrope"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = font
    return tb


def add_bullets(slide, left, top, w, h, items, size=16, color=TEXT):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Manrope"
    return tb


def add_page_num(slide, n):
    add_text(slide, Inches(12.3), Inches(7.0), Inches(1), Inches(0.4),
             str(n), size=10, color=FAINT, align=PP_ALIGN.RIGHT)


def add_title_bar(slide, title, subtitle=None):
    """每页顶部标题栏"""
    add_rect(slide, Inches(0.6), Inches(0.4), Inches(0.08), Inches(0.5), ACCENT)
    add_text(slide, Inches(0.85), Inches(0.35), Inches(10), Inches(0.6),
             title, size=26, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.85), Inches(0.85), Inches(10), Inches(0.4),
                 subtitle, size=13, color=DIM)


# ========== 第1页: 封面 ==========
s = add_slide()
add_rect(s, Inches(0), Inches(0), SW, SH, BG)
# 装饰球
from pptx.enum.shapes import MSO_SHAPE
orb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.5), Inches(5), Inches(5))
orb.fill.solid()
orb.fill.fore_color.rgb = RGBColor(0x0A, 0x3D, 0x62)
orb.line.color.rgb = ACCENT
orb.line.width = Pt(1.5)
orb.shadow.inherit = False
add_text(s, Inches(0.8), Inches(2.2), Inches(7), Inches(1.2),
         "🧠 魔法大脑", size=54, color=ACCENT, bold=True)
add_text(s, Inches(0.8), Inches(3.4), Inches(7), Inches(0.8),
         "Magic Brain", size=24, color=DIM)
add_text(s, Inches(0.8), Inches(4.2), Inches(7), Inches(0.6),
         "AI 交通套利搜索引擎", size=22, color=TEXT, bold=True)
add_text(s, Inches(0.8), Inches(4.9), Inches(7), Inches(1),
         "输入起终点，自动发现更便宜的中转路线\n上海→香港 省 ¥399 (41%)", size=15, color=SAVE)
add_text(s, Inches(0.8), Inches(6.3), Inches(7), Inches(0.5),
         "GOAI 无界应用｜Boundless Agents · 出行规划细分场景", size=12, color=FAINT)

# ========== 第2页: 故事 ==========
s = add_slide()
add_title_bar(s, "一个真实的省钱痛点", "上海 → 香港")
# 三方案对比卡片
y = Inches(1.8)
cards = [
    ("🚄 高铁直达 G99", "¥966", "8.5小时 · 直达", False),
    ("✈️ 最便宜航班", "¥696", "3小时 · 直达", False),
    ("🛏️ 魔法大脑方案", "¥567", "动卧经深圳过关", True),
]
for i, (name, price, meta, hl) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    w = Inches(3.7)
    color = RGBColor(0x16, 0xA3, 0x4A) if hl else CARD
    line = SAVE if hl else None
    add_rect(s, x, y, w, Inches(3), color, line=line)
    add_text(s, x, Inches(2.1), w, Inches(0.5), name, size=15, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.8), w, Inches(1), price, size=44 if hl else 36, color=SAVE if hl else ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(4.0), w, Inches(0.5), meta, size=12, color=DIM, align=PP_ALIGN.CENTER)
    if hl:
        add_text(s, x, Inches(4.5), w, Inches(0.4), "⭐ 推荐", size=13, color=SAVE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.8),
         "省 ¥399 (41%)  ·  还能在卧铺睡一觉省一晚酒店", size=20, color=SAVE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
         "这不是个例——多模态中转里藏着大量被忽视的省钱机会", size=14, color=DIM, align=PP_ALIGN.CENTER)
add_page_num(s, 2)

# ========== 第3页: 市场空白 ==========
s = add_slide()
add_title_bar(s, "市场空白——为什么没人做这件事？")
rows = [
    ("导航 App", "高德/百度/Google", "只做最快路径", "不做省钱、不做跨交通方式组合"),
    ("OTA 平台", "携程/飞猪/12306", "卖直达库存", "不主动推荐绕路省钱"),
    ("机票比价", "天巡/去哪儿", "单一航班比价", "不组合火车+飞机+过关"),
]
y = Inches(1.8)
add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.55), CARD)
for i, h in enumerate(["产品类型", "代表", "它做什么", "它不做"]):
    add_text(s, Inches(0.9 + i*2.9), Inches(1.9), Inches(2.8), Inches(0.4), h, size=14, color=ACCENT, bold=True)
for r, row in enumerate(rows):
    ry = Inches(2.5 + r*0.9)
    add_rect(s, Inches(0.8), ry, Inches(11.7), Inches(0.8), CARD, line=RGBColor(0x1E,0x29,0x3B))
    for i, cell in enumerate(row):
        c = TEXT if i < 2 else (SAVE if i == 2 else RGBColor(0xF8,0x71,0x71))
        add_text(s, Inches(0.9 + i*2.9), Inches(2.65 + r*0.9), Inches(2.8), Inches(0.5), cell, size=13, color=c)
add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8),
         "结论：没有产品做「多模态中转套利」——这是魔法大脑的位置", size=18, color=SAVE, bold=True, align=PP_ALIGN.CENTER)
add_page_num(s, 3)

# ========== 第4页: Agent 闭环 ==========
s = add_slide()
add_title_bar(s, "4 Agent 协作的任务闭环", "从自然语言到省钱方案")
agents = [
    ("🧠", "意图理解 Agent", "GLM 解析起终点、偏好、约束"),
    ("📋", "任务规划 Agent", "GLM 拆解：查直达 → 找中转 → 对比"),
    ("🔧", "工具调用 Agent", "DuckDuckGo 联网查票价 + 本地多目标搜索"),
    ("💡", "结果解释 Agent", "GLM 生成自然语言回复 + 路线卡"),
]
for i, (icon, name, desc) in enumerate(agents):
    y = Inches(1.8 + i*1.15)
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(1), CARD, line=RGBColor(0x1E,0x29,0x3B))
    add_text(s, Inches(1.0), Inches(2.0 + i*1.15), Inches(0.6), Inches(0.6), icon, size=24, color=TEXT)
    add_text(s, Inches(1.8), Inches(1.95 + i*1.15), Inches(3), Inches(0.5), name, size=16, color=ACCENT, bold=True)
    add_text(s, Inches(1.8), Inches(2.35 + i*1.15), Inches(10), Inches(0.4), desc, size=13, color=DIM)
    if i < 3:
        add_text(s, Inches(6), Inches(2.75 + i*1.15), Inches(1), Inches(0.3), "↓", size=16, color=ACCENT, align=PP_ALIGN.CENTER)
add_page_num(s, 4)

# ========== 第5页: Agent 能力清单 ==========
s = add_slide()
add_title_bar(s, "Agent 能力——完整任务闭环")
caps = [
    ("任务输入", "自然语言（下周三晚下班去香港）或结构化（起终点下拉）"),
    ("意图理解", "GLM-4.7 解析出发地、目的地、偏好、时间约束"),
    ("任务规划", "GLM 拆解搜索步骤，说明需调用哪些工具"),
    ("工具调用", "DuckDuckGo 联网查实时票价 + 本地多目标 Dijkstra"),
    ("结果交付", "省钱横幅 + 多方案对比 + 路线卡 + 3D 地球路径"),
    ("验证反馈", "价格来源可追溯、核实状态标注、参考性提示"),
    ("安全边界", "不订票、不替代 OTA、联网结果标注以购票平台为准"),
]
for i, (k, v) in enumerate(caps):
    y = Inches(1.7 + i*0.72)
    add_rect(s, Inches(0.8), y, Inches(2.2), Inches(0.55), RGBColor(0x0C,0x1A,0x2E), line=ACCENT)
    add_text(s, Inches(0.9), Inches(1.78 + i*0.72), Inches(2), Inches(0.4), k, size=13, color=ACCENT, bold=True)
    add_text(s, Inches(3.2), Inches(1.78 + i*0.72), Inches(9.5), Inches(0.4), v, size=13, color=TEXT)
add_page_num(s, 5)

# ========== 第6页: 核心算法 ==========
s = add_slide()
add_title_bar(s, "核心算法——多模态交通套利搜索")
add_text(s, Inches(0.8), Inches(1.7), Inches(6), Inches(0.5), "问题建模：多模态交通网络", size=18, color=ACCENT, bold=True)
add_bullets(s, Inches(0.8), Inches(2.3), Inches(6), Inches(2.5), [
    "节点 = 城市/枢纽/口岸站点",
    "边 = 一段交通（飞机/高铁/卧铺/地铁/过关）",
    "每条边带多维成本：价格、耗时、换乘、舒适度",
], size=15)
add_text(s, Inches(0.8), Inches(4.2), Inches(6), Inches(0.5), "算法：多目标 Dijkstra 变种", size=18, color=ACCENT, bold=True)
add_bullets(s, Inches(0.8), Inches(4.8), Inches(6), Inches(2.2), [
    "综合成本 = w₁·价格 + w₂·时间 + w₃·换乘 + w₄·舒适度",
    "用户偏好（省钱/时间/均衡）调整权重",
    "输出 Top-K 条 Pareto 候选，与直达基准对比",
], size=15)
# 右侧公式卡片
add_rect(s, Inches(7.5), Inches(2.0), Inches(5), Inches(3.8), CARD, line=ACCENT)
add_text(s, Inches(7.7), Inches(2.2), Inches(4.6), Inches(0.5), "关键差异", size=16, color=SAVE, bold=True)
add_text(s, Inches(7.7), Inches(2.8), Inches(4.6), Inches(2),
         "不是最短路\n\n而是「省钱空间最大」的路\n\n普通导航找最快\nOTA 只卖直达\n魔法大脑找套利组合", size=15, color=TEXT)
add_page_num(s, 6)

# ========== 第7页: 3D 地球 (带截图) ==========
s = add_slide()
add_title_bar(s, "3D 地球真实路径可视化", "多模态跨城出行 = 地理问题")
if (SHOTS / "01-home.png").exists():
    s.shapes.add_picture(str(SHOTS / "01-home.png"), Inches(0.8), Inches(1.7), width=Inches(7.5))
add_rect(s, Inches(8.6), Inches(1.7), Inches(4), Inches(4.2), CARD, line=ACCENT)
add_text(s, Inches(8.8), Inches(1.9), Inches(3.6), Inches(0.5), "为什么是地球？", size=16, color=ACCENT, bold=True)
add_bullets(s, Inches(8.8), Inches(2.5), Inches(3.6), Inches(3.5), [
    "Three.js 真实地球（NASA Blue Marble）",
    "火车贴地：沿真实城市航点",
    "飞机大圆弧：起终点接地",
    "途径站点可悬停显示",
    "省钱横幅醒目（¥399 / 41%↓）",
], size=13)
add_page_num(s, 7)

# ========== 第8页: 技术架构 ==========
s = add_slide()
add_title_bar(s, "技术架构")
layers = [
    ("前端层", "Three.js + Tailwind · 3D 地球 · 毛玻璃 UI · Agent 思考可视化", ACCENT),
    ("服务层", "FastAPI · /api/agent · /api/search · /api/geo", ACCENT2),
    ("Agent 层", "4 Agent pipeline + GLM-4.7 (Anthropic 兼容) + DuckDuckGo 联网", SAVE),
    ("引擎层", "交通图 · 多目标 Dijkstra · 套利分析（纯标准库）", ACCENT),
    ("数据层", "JSON 线路数据 + 89 城市经纬度", ACCENT2),
]
for i, (name, desc, color) in enumerate(layers):
    y = Inches(1.8 + i*1.0)
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.85), CARD, line=color)
    add_rect(s, Inches(0.8), y, Inches(1.8), Inches(0.85), color)
    add_text(s, Inches(0.9), Inches(2.0 + i*1.0), Inches(1.6), Inches(0.5), name, size=14, color=BG, bold=True)
    add_text(s, Inches(2.8), Inches(2.0 + i*1.0), Inches(9.5), Inches(0.5), desc, size=13, color=TEXT)
add_page_num(s, 8)

# ========== 第9页: 真实数据 (带截图) ==========
s = add_slide()
add_title_bar(s, "真实数据与可追溯")
if (SHOTS / "03-agent.png").exists():
    s.shapes.add_picture(str(SHOTS / "03-agent.png"), Inches(0.8), Inches(1.7), width=Inches(7.2))
add_rect(s, Inches(8.3), Inches(1.7), Inches(4.3), Inches(4.5), CARD, line=ACCENT)
add_text(s, Inches(8.5), Inches(1.9), Inches(3.9), Inches(0.5), "数据来源", size=16, color=ACCENT, bold=True)
add_bullets(s, Inches(8.5), Inches(2.5), Inches(3.9), Inches(3.8), [
    "本地核实：上海→香港（12306 实测 ✅）",
    "联网实时：DuckDuckGo 查公开网页",
    "来源含网页标题与链接，可追溯",
    "任意城市可搜（联网模式）",
    "待核实线路标注 ⏳",
], size=12)
add_page_num(s, 9)

# ========== 第10页: 合规 ==========
s = add_slide()
add_title_bar(s, "安全合规与边界")
items = [
    ("定位", "出行规划辅助工具，非订票平台"),
    ("不替代专业决策", "不订票、不出票、不替代 OTA/12306 最终票价"),
    ("风险提示", "联网结果标注「以购票平台为准」"),
    ("隐私", "无账号、无数据存储、无追踪"),
    ("第三方披露", "GLM-4.7（可迁移）、Three.js/FastAPI（MIT）"),
]
for i, (k, v) in enumerate(items):
    y = Inches(1.8 + i*0.95)
    add_rect(s, Inches(0.8), y, Inches(2.5), Inches(0.75), RGBColor(0x0C,0x1A,0x2E), line=SAVE)
    add_text(s, Inches(0.9), Inches(2.0 + i*0.95), Inches(2.3), Inches(0.4), k, size=13, color=SAVE, bold=True)
    add_text(s, Inches(3.5), Inches(2.0 + i*0.95), Inches(9), Inches(0.4), v, size=14, color=TEXT)
add_page_num(s, 10)

# ========== 第11页: 开放复用 ==========
s = add_slide()
add_title_bar(s, "开放 / 复用贡献")
items = [
    ("多模态交通图引擎", "纯标准库，可独立用于任意路径规划场景"),
    ("4 Agent 流水线架构", "意图→规划→工具→解释，可复用于其他工具调用场景"),
    ("真实交通数据集", "核实线路 + 89 城市经纬度，可作示例数据"),
    ("GLM Agent 构建范例", "代码即文档，展示如何用国产开源模型构建联网 Agent"),
]
for i, (k, v) in enumerate(items):
    y = Inches(1.9 + i*1.15)
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(1), CARD, line=ACCENT)
    add_text(s, Inches(1.0), Inches(2.1 + i*1.15), Inches(0.6), Inches(0.6), "📦", size=20)
    add_text(s, Inches(1.8), Inches(2.0 + i*1.15), Inches(4), Inches(0.5), k, size=16, color=ACCENT, bold=True)
    add_text(s, Inches(1.8), Inches(2.4 + i*1.15), Inches(10), Inches(0.4), v, size=13, color=DIM)
add_text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
         "MIT License · GitHub 公开", size=14, color=SAVE, bold=True, align=PP_ALIGN.CENTER)
add_page_num(s, 11)

# ========== 第12页: Demo ==========
s = add_slide()
add_title_bar(s, "Demo 演示", "完整任务闭环")
demos = [
    "搜索框输入上海→香港 → Agent 4 步思考 → 省 ¥399",
    "点方案卡 → 地球画真实路径（火车贴地）",
    "切时间优先 → 横幅切换为省时间",
    "AI 对话：「下周三晚下班去香港」→ Agent 联网 → 路线卡",
    "换小众城市（成都→西安）→ 联网查到真实票价",
]
for i, d in enumerate(demos):
    y = Inches(1.9 + i*0.85)
    add_rect(s, Inches(0.8), y, Inches(0.6), Inches(0.6), ACCENT)
    add_text(s, Inches(0.85), Inches(1.98 + i*0.85), Inches(0.5), Inches(0.5), str(i+1), size=18, color=BG, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), Inches(2.0 + i*0.85), Inches(11), Inches(0.5), d, size=15, color=TEXT)
add_page_num(s, 12)

# ========== 第13页: 后续计划 ==========
s = add_slide()
add_title_bar(s, "后续计划")
plans = [
    ("📊 数据扩充", "更多核实线路 + Agent 定时联网更新票价"),
    ("⚡ 实时性", "接入真实 12306/航司 API（如有授权）"),
    ("📱 移动端", "响应式适配 + PWA"),
    ("👤 个性化", "用户偏好画像（出行习惯学习）"),
    ("🌐 开放", "社区众包线路、开放 API 供第三方接入"),
]
for i, (k, v) in enumerate(plans):
    y = Inches(1.9 + i*0.95)
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.8), CARD, line=RGBColor(0x1E,0x29,0x3B))
    add_text(s, Inches(1.0), Inches(2.05 + i*0.95), Inches(3), Inches(0.5), k, size=15, color=ACCENT, bold=True)
    add_text(s, Inches(4.2), Inches(2.05 + i*0.95), Inches(8), Inches(0.5), v, size=14, color=TEXT)
add_page_num(s, 13)

# ========== 第14页: 致谢 ==========
s = add_slide()
add_text(s, Inches(2), Inches(2.8), Inches(9), Inches(1), "🧠 魔法大脑", size=48, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(4.0), Inches(9), Inches(0.6), "让每一程更省", size=22, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(5.2), Inches(9), Inches(0.5), "感谢评审 · 期待反馈", size=16, color=DIM, align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print(f"PPT 已生成: {OUT}")
print(f"共 {len(prs.slides)} 页")
